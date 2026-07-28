"""Office document extraction (docx, xlsx, pptx)."""

from pathlib import Path


def extract_docx(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception:
        return ""


def extract_xlsx(path: str) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, data_only=True, read_only=True)
        lines: list[str] = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                text = " ".join(str(v) for v in row if v is not None)
                if text.strip():
                    lines.append(text)
        return "\n".join(lines)
    except Exception:
        return ""


def extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            lines.append(t)
        return "\n".join(lines)
    except Exception:
        return ""


SUPPORTED: dict[str, callable] = {
    ".docx": extract_docx,
    ".doc": extract_docx,
    ".xlsx": extract_xlsx,
    ".pptx": extract_pptx,
}


def extract(path: str) -> str:
    ext = Path(path).suffix.lower()
    handler = SUPPORTED.get(ext)
    if handler:
        return handler(path)
    return ""
